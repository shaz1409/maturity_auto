#!/usr/bin/env python3
"""
Marketing Maturity Assessment Automation

Automated script to:
1. Load survey data from Google Sheets
2. Calculate category scores for each client
3. Generate AI-powered recommendations
4. Create personalized PowerPoint presentations

Usage:
    python maturity_assessment.py

Environment Variables:
    OPENAI_API_KEY: Your OpenAI API key (required)
"""

import pandas as pd
import numpy as np
import requests
import io
import urllib3
import os
import re
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt

# SharePoint imports (optional)
try:
    from office365.sharepoint.client_context import ClientContext
    from office365.runtime.auth.user_credential import UserCredential
    SHAREPOINT_AVAILABLE = True
except ImportError:
    try:
        # Fallback to older import
        from office365.sharepoint.client_context import ClientContext
        from office365.runtime.auth.authentication_context import AuthenticationContext
        SHAREPOINT_AVAILABLE = True
        USE_LEGACY_AUTH = True
    except ImportError:
        SHAREPOINT_AVAILABLE = False
        USE_LEGACY_AUTH = False
else:
    USE_LEGACY_AUTH = False

# Suppress SSL warnings
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# Configuration
SHEET_ID = '1tHWUJWJl_zTwGRTg21qbW_5qpYH8bBMFYZYbIZ03eO8'
GID = '491555971'
SHEET_URL = f'https://docs.google.com/spreadsheets/d/{SHEET_ID}/export?format=csv&gid={GID}'
TEMPLATE_PATH = 'Maturity_Slide_Template.pptx'
OUTPUT_DIR = 'output'

# SharePoint configuration
SHAREPOINT_SITE_URL = 'https://indigitalmarketing.sharepoint.com/sites/team'
SHAREPOINT_FOLDER_PATH = '/Drive/Team/Shaz/Client_Maturity_Automation'
SHAREPOINT_UPLOAD_ENABLED = os.environ.get('SHAREPOINT_UPLOAD', 'false').lower() == 'true'

# SharePoint authentication method: 'user' (username/password) or 'app' (client ID/secret)
SHAREPOINT_AUTH_METHOD = os.environ.get('SHAREPOINT_AUTH_METHOD', 'user').lower()

# Initialize OpenAI client
try:
    from openai import OpenAI
    api_key = os.environ.get('OPENAI_API_KEY')
    if not api_key or api_key == 'your-api-key-here':
        print("ERROR: OPENAI_API_KEY not set. Please set it as an environment variable.")
        sys.exit(1)
    client = OpenAI(api_key=api_key)
except ImportError:
    print("ERROR: OpenAI package not installed. Run: pip install openai")
    sys.exit(1)
except Exception as e:
    print(f"ERROR: Failed to initialize OpenAI client: {e}")
    sys.exit(1)


def clean_column_name(col_name):
    """Clean column name: lowercase, remove ALL special characters, replace spaces with underscores"""
    cleaned = col_name.lower()
    cleaned = cleaned.replace('-', ' ')
    cleaned = re.sub(r'[^a-z0-9\s]', '', cleaned)
    cleaned = re.sub(r'\s+', ' ', cleaned)
    cleaned = cleaned.strip().replace(' ', '_')
    return cleaned


def load_data():
    """Load survey data from Google Sheets"""
    print("Loading data from Google Sheet...")
    response = requests.get(SHEET_URL, verify=False, timeout=30)
    response.raise_for_status()
    df = pd.read_csv(io.StringIO(response.text))
    print(f"✓ Data loaded: {df.shape[0]} rows, {df.shape[1]} columns")
    return df


def setup_mappings(df):
    """Set up column name mappings and question categories"""
    # Create column name mappings
    COLUMN_NAME_MAPPING = {}
    CLEANED_TO_ORIGINAL = {}

    for col in df.columns:
        cleaned = clean_column_name(col)
        COLUMN_NAME_MAPPING[col] = cleaned
        CLEANED_TO_ORIGINAL[cleaned] = col

    # Get question columns in order
    question_columns = [col for col in df.columns if col not in ['Timestamp', 'Email Address']]
    cleaned_questions = [COLUMN_NAME_MAPPING[q] for q in question_columns]

    # Define question categories
    QUESTION_CATEGORIES = {
        'Tech & Data': cleaned_questions[0:5],
        'Campaigning & Assets': cleaned_questions[5:11],
        'Segmentation & Personalisation': cleaned_questions[11:14],
        'Reporting & Insights': cleaned_questions[14:20],
        'People & Operations': cleaned_questions[20:24]
    }

    # Create reverse mapping
    QUESTION_TO_CATEGORY = {}
    for category, questions in QUESTION_CATEGORIES.items():
        for question in questions:
            QUESTION_TO_CATEGORY[question] = category

    # Create cleaned to original column mapping
    CLEANED_TO_ORIGINAL_COL = {cleaned: original for original, cleaned in COLUMN_NAME_MAPPING.items()}

    return (COLUMN_NAME_MAPPING, CLEANED_TO_ORIGINAL_COL, QUESTION_CATEGORIES, 
            QUESTION_TO_CATEGORY, question_columns)


def calculate_category_scores(client_row, question_categories, cleaned_to_original):
    """Calculate average score for each category for a single client"""
    scores = {}

    for category, questions in question_categories.items():
        category_scores = []

        for cleaned_q in questions:
            original_col = cleaned_to_original.get(cleaned_q)

            if original_col:
                value = client_row.get(original_col, None)

                if value is not None and pd.notna(value):
                    try:
                        num_value = float(value)
                        if 1 <= num_value <= 4:
                            category_scores.append(num_value)
                    except (ValueError, TypeError):
                        pass

        if category_scores:
            scores[category] = np.mean(category_scores)
        else:
            scores[category] = None

    return scores


def find_text_boxes(slide):
    """Find text boxes and orange circle on a slide"""
    elements = {
        'score': None,
        'recommendations': None,
        'orange_circle': None,
        'line': None
    }

    # Find text boxes
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if "Your score" in text and elements['score'] is None:
                elements['score'] = shape
            elif "Recommendation 1" in text and elements['recommendations'] is None:
                elements['recommendations'] = shape

    # Find line and circle
    lines = []
    circles = []

    for shape in slide.shapes:
        if hasattr(shape, 'shape_type'):
            if shape.shape_type == 9:  # LINE
                lines.append(shape)
            elif shape.shape_type == 1:  # AUTO_SHAPE
                if hasattr(shape, 'width') and hasattr(shape, 'height'):
                    width = shape.width
                    height = shape.height
                    if abs(width - height) / max(width, height) < 0.2:
                        try:
                            if hasattr(shape, 'fill') and shape.fill.type == 1:
                                circles.append(shape)
                        except:
                            pass

    if lines:
        elements['line'] = max(lines, key=lambda l: l.width if hasattr(l, 'width') else 0)

    if circles and elements['line']:
        line_y = elements['line'].top + (elements['line'].height / 2)
        elements['orange_circle'] = min(circles, 
                                        key=lambda c: abs((c.top + c.height/2) - line_y))
    elif circles:
        elements['orange_circle'] = circles[0]

    return elements


def generate_recommendations(category, score, questions_in_category, client_responses, original_questions_dict):
    """Generate recommendations using OpenAI"""
    # Identify low-scoring questions
    low_scoring_questions = []
    all_questions_context = []

    for cleaned_q in questions_in_category:
        q_score = client_responses.get(cleaned_q, None)
        if q_score is not None and isinstance(q_score, (int, float)):
            original_q = original_questions_dict.get(cleaned_q, cleaned_q)
            question_info = {
                'cleaned': cleaned_q,
                'original': original_q,
                'score': q_score
            }
            all_questions_context.append(question_info)

            if q_score <= 2:
                low_scoring_questions.append(question_info)

    low_scoring_questions.sort(key=lambda x: x['score'])
    all_questions_context.sort(key=lambda x: x['score'])

    questions_detail = "\n".join([
        f"- Question: {q['original']}\n  Score: {q['score']}/4"
        for q in all_questions_context
    ])

    focus_areas = ""
    if low_scoring_questions:
        focus_areas = "\n\nAreas requiring immediate attention (low scores):\n" + "\n".join([
            f"- {q['original']} (Score: {q['score']}/4)"
            for q in low_scoring_questions[:3]
        ])

    # Determine maturity level
    if score <= 1.5:
        maturity_level = "not mature"
    elif score <= 2.5:
        maturity_level = "developing"
    elif score <= 3.5:
        maturity_level = "mature"
    else:
        maturity_level = "very mature"

    prompt = f"""You are a CRM marketing maturity consultant. Generate recommendations for a client.

Category: {category}
Overall Maturity Score: {score:.2f}/4.0 ({maturity_level})

Client's responses to all questions in this category:
{questions_detail}
{focus_areas}

Instructions:
1. Generate a brief 2-3 sentence summary of their current maturity level in this category
2. Generate four specific, actionable recommendations that:
   - PRIORITIZE addressing the low-scoring areas identified above
   - Reference the specific questions where they scored low (1-2 out of 4)
   - Provide concrete, actionable steps based on the question context
   - Are tailored to their current maturity level

Focus especially on the questions where they scored 1-2, as these are the areas needing the most improvement.

Format the response as:
SUMMARY: [your summary here]
RECOMMENDATIONS:
1. [recommendation 1 - should address a specific low-scoring question]
2. [recommendation 2 - should address a specific low-scoring question]
3. [recommendation 3 - can address another area or build on improvements]
4. [recommendation 4 - can address another area or build on improvements]

Make each recommendation specific, actionable, and directly related to the questions they answered poorly."""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are an expert CRM marketing consultant providing actionable recommendations."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=500
        )

        result = response.choices[0].message.content

        if "SUMMARY:" in result:
            parts = result.split("RECOMMENDATIONS:")
            summary = parts[0].replace("SUMMARY:", "").strip()
            recommendations_text = parts[1] if len(parts) > 1 else ""

            recommendations = []
            for line in recommendations_text.split('\n'):
                line = line.strip()
                if line and (line[0].isdigit() or line.startswith('-')):
                    rec = line.split('.', 1)[-1].strip() if '.' in line else line.lstrip('- ').strip()
                    if rec:
                        recommendations.append(rec)

            while len(recommendations) < 4:
                recommendations.append("Continue building on the recommendations above.")

            return summary[:200], recommendations[:4]
        else:
            lines = result.split('\n')
            summary = lines[0][:200] if lines else "Summary generated"
            recommendations = [line.strip() for line in lines[1:] if line.strip() and len(line.strip()) > 10][:4]
            while len(recommendations) < 4:
                recommendations.append("Continue improving in this area.")
            return summary, recommendations

    except Exception as e:
        print(f"  ⚠️  Error generating recommendations: {e}")
        return f"Error: {str(e)}", [
            "Recommendation 1: Review current processes",
            "Recommendation 2: Identify improvement areas",
            "Recommendation 3: Implement best practices",
            "Recommendation 4: Monitor progress"
        ]


def map_slides_to_categories(prs):
    """Map slide titles to category names"""
    SLIDE_CATEGORY_MAPPING = {
        'Tech and Data': 'Tech & Data',
        'Campaigning & Assets': 'Campaigning & Assets',
        'Segmentation & Personalisation': 'Segmentation & Personalisation',
        'Reporting & Insights': 'Reporting & Insights',
        'People & Operations': 'People & Operations'
    }
    
    CATEGORY_TO_SLIDE = {}
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                title = shape.text.strip()
                if title in SLIDE_CATEGORY_MAPPING:
                    category = SLIDE_CATEGORY_MAPPING[title]
                    CATEGORY_TO_SLIDE[category] = i
                    break
    
    return CATEGORY_TO_SLIDE


def email_to_filename(email):
    """Convert email to safe filename"""
    safe_email = email.replace('@', '_at_').replace('.', '_')
    return f"{safe_email}_Maturity_Assessment.pptx"


def upload_to_sharepoint(file_path, filename):
    """
    Upload a file to SharePoint folder (only if it doesn't already exist)

    Parameters:
    - file_path: Local path to the file
    - filename: Name for the file in SharePoint
    
    Returns:
    - True if successful, False otherwise
    """
    if not SHAREPOINT_UPLOAD_ENABLED:
        return False
    
    if not SHAREPOINT_AVAILABLE:
        print(f"    ⚠️  SharePoint library not available. Install with: pip install Office365-REST-Python-Client")
        return False
    
    try:
        # Authenticate based on method
        print(f"      Authenticating to SharePoint...")
        
        if SHAREPOINT_AUTH_METHOD == 'app':
            # App-based authentication (recommended for modern SharePoint)
            client_id = os.environ.get('SHAREPOINT_CLIENT_ID')
            client_secret = os.environ.get('SHAREPOINT_CLIENT_SECRET')
            tenant_id = os.environ.get('SHAREPOINT_TENANT_ID', 'common')
            
            if not client_id or not client_secret:
                print(f"    ⚠️  App credentials not set (SHAREPOINT_CLIENT_ID, SHAREPOINT_CLIENT_SECRET)")
                print(f"    Set SHAREPOINT_AUTH_METHOD=app and provide app credentials")
                return False
            
            try:
                from office365.runtime.auth.client_credential import ClientCredential
                credentials = ClientCredential(client_id, client_secret)
                ctx = ClientContext(SHAREPOINT_SITE_URL).with_credentials(credentials)
                
                # Test connection
                print(f"      Testing connection...")
                web = ctx.web
                ctx.load(web)
                ctx.execute_query()
                print(f"      ✓ App-based authentication successful")
            except Exception as app_error:
                error_msg = str(app_error)
                print(f"      ❌ App-based authentication failed: {app_error}")
                
                # Provide helpful guidance for common errors
                if "403" in error_msg or "forbidden" in error_msg.lower():
                    print(f"      ⚠️  403 Forbidden - Permission issue:")
                    print(f"         - Make sure you're using APPLICATION permissions (not Delegated)")
                    print(f"         - Grant 'Sites.ReadWrite.All' as APPLICATION permission")
                    print(f"         - Admin consent must be granted")
                    print(f"         - The app needs access to the SharePoint site")
                elif "401" in error_msg or "unauthorized" in error_msg.lower():
                    print(f"      ⚠️  Check your Client ID and Client Secret")
                elif "invalid" in error_msg.lower() or "client" in error_msg.lower():
                    print(f"      ⚠️  Invalid credentials - verify Client ID and Secret")
                
                return False
                
        else:
            # Username/password authentication
            username = os.environ.get('SHAREPOINT_USERNAME')
            password = os.environ.get('SHAREPOINT_PASSWORD')
            
            if not username or not password:
                print(f"    ⚠️  SharePoint credentials not set (SHAREPOINT_USERNAME, SHAREPOINT_PASSWORD)")
                return False
            
            try:
                from office365.runtime.auth.user_credential import UserCredential
                credentials = UserCredential(username, password)
                ctx = ClientContext(SHAREPOINT_SITE_URL).with_credentials(credentials)
                
                # Test the connection
                print(f"      Testing connection...")
                web = ctx.web
                ctx.load(web)
                ctx.execute_query()
                print(f"      ✓ User authentication successful")
                
            except Exception as auth_error:
                error_msg = str(auth_error)
                print(f"      ❌ SharePoint authentication/connection failed")
                print(f"      Error: {error_msg}")
                
                # Provide helpful guidance
                if "binary security token" in error_msg.lower() or "extSTS" in error_msg:
                    print(f"      ⚠️  This error means legacy auth is not supported.")
                    print(f"      Solution: Use app-based authentication instead:")
                    print(f"        1. Register an app in Azure AD")
                    print(f"        2. Set SHAREPOINT_AUTH_METHOD=app")
                    print(f"        3. Set SHAREPOINT_CLIENT_ID and SHAREPOINT_CLIENT_SECRET")
                elif "401" in error_msg or "unauthorized" in error_msg.lower():
                    print(f"      ⚠️  Check your username and password")
                elif "403" in error_msg or "forbidden" in error_msg.lower():
                    print(f"      ⚠️  Account may not have permission to access this site")
                
                return False
        
    except Exception as e:
        print(f"    ⚠️  SharePoint setup error: {e}")
        return False
    
    # Get the folder
    folder = ctx.web.get_folder_by_server_relative_url(SHAREPOINT_FOLDER_PATH)
    
    # Check if file already exists in SharePoint
    try:
        print(f"      Checking if file exists in SharePoint...")
        files = folder.files
        ctx.load(files)
        ctx.execute_query()
        
        # Check if file with this name already exists
        file_exists = any(f.properties['Name'] == filename for f in files)
        
        if file_exists:
            print(f"      ⊘ File already exists in SharePoint - skipping upload")
            return True  # Return True since file is already there
        
        print(f"      ✓ File not found in SharePoint - will upload")
        
    except Exception as check_error:
        # If we can't check, proceed with upload anyway
        print(f"      ⚠️  Could not check if file exists: {check_error}")
        print(f"      → Proceeding with upload anyway")
    
    # Read file content
    try:
        print(f"      Reading file...")
        file_size = os.path.getsize(file_path)
        print(f"      File size: {file_size / 1024:.1f} KB")
        
        with open(file_path, 'rb') as file_content:
            file_data = file_content.read()
        
        # Upload file
        print(f"      Uploading to SharePoint...")
        uploaded_file = folder.upload_file(filename, file_data).execute_query()
        
        print(f"      ✓ Successfully uploaded to SharePoint: {filename}")
        return True
        
    except Exception as e:
        print(f"    ⚠️  SharePoint upload error: {e}")
        return False


def generate_client_presentation(client_email, client_scores, client_responses, 
                                  COLUMN_NAME_MAPPING, CLEANED_TO_ORIGINAL_COL,
                                  QUESTION_CATEGORIES, output_filename):
    """Generate PowerPoint presentation for a client"""
    prs = Presentation(TEMPLATE_PATH)
    CATEGORY_TO_SLIDE = map_slides_to_categories(prs)
    
    # Map responses
    question_responses = {}
    for orig_col, cleaned_col in COLUMN_NAME_MAPPING.items():
        if orig_col in client_responses:
            question_responses[cleaned_col] = client_responses[orig_col]

    # Process each category
    for category, score in client_scores.items():
        if category not in CATEGORY_TO_SLIDE:
            continue

        if pd.isna(score) or score is None:
            continue

        slide_idx = CATEGORY_TO_SLIDE[category]
        slide = prs.slides[slide_idx]
        elements = find_text_boxes(slide)

        questions_in_category = QUESTION_CATEGORIES.get(category, [])
        category_responses = {q: question_responses.get(q, 'N/A') for q in questions_in_category}

        original_questions_dict = {}
        for cleaned_q in questions_in_category:
            orig_col = CLEANED_TO_ORIGINAL_COL.get(cleaned_q)
            if orig_col:
                original_questions_dict[cleaned_q] = orig_col

        # Generate recommendations
        print(f"  Generating recommendations for {category} (score: {score:.2f})...")
        summary, recommendations = generate_recommendations(
            category, score, questions_in_category, category_responses, original_questions_dict
        )
        
        # Update score
        if elements['score']:
            if hasattr(elements['score'], 'text_frame'):
                elements['score'].text_frame.clear()
                p = elements['score'].text_frame.paragraphs[0]
                p.text = f"{score:.2f}/4.0"
            else:
                elements['score'].text = f"{score:.2f}/4.0"

        # Update recommendations with smaller font to fit on slide
        if elements['recommendations']:
            recommendations_text = "\n\n".join([f"{i+1}. {rec}" for i, rec in enumerate(recommendations)])
            if hasattr(elements['recommendations'], 'text_frame'):
                text_frame = elements['recommendations'].text_frame
                text_frame.clear()
                
                # Enable word wrap
                text_frame.word_wrap = True
                
                # Split by double newlines to create paragraphs for each recommendation
                paragraphs = recommendations_text.split("\n\n")
                
                for i, para_text in enumerate(paragraphs):
                    if i > 0:
                        p = text_frame.add_paragraph()
                    else:
                        p = text_frame.paragraphs[0]
                    
                    p.text = para_text
                    p.space_after = Pt(3)  # Small spacing between recommendations
                    
                    # Set smaller font size (9pt) for all runs in this paragraph
                    for run in p.runs:
                        run.font.size = Pt(9)
                
            else:
                elements['recommendations'].text = recommendations_text

        # Position orange circle
        if elements['orange_circle'] and elements['line']:
            try:
                line = elements['line']
                line_left = line.left
                line_width = line.width
                line_top = line.top
                line_height = line.height

                score_position = score / 4.0
                circle_x = line_left + (line_width * score_position)
                circle_y = line_top + (line_height / 2)

                circle = elements['orange_circle']
                circle_width = circle.width
                circle_height = circle.height

                circle.left = int(circle_x - (circle_width / 2))
                circle.top = int(circle_y - (circle_height / 2))
            except Exception as e:
                print(f"    ⚠️  Could not position orange circle: {e}")

    prs.save(output_filename)
    return output_filename


def main():
    """Main execution function"""
    print("="*60)
    print("Marketing Maturity Assessment Automation")
    print(f"Started at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
    print("="*60)
    
    # Create output directory
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    try:
        # Load data
        df = load_data()
        
        # Setup mappings
        print("\nSetting up mappings...")
        (COLUMN_NAME_MAPPING, CLEANED_TO_ORIGINAL_COL, QUESTION_CATEGORIES,
         QUESTION_TO_CATEGORY, question_columns) = setup_mappings(df)
        print(f"✓ Mapped {len(QUESTION_CATEGORIES)} categories")
        
        # Calculate scores for all clients
        print("\nCalculating category scores...")
        category_scores_list = []
        
        for idx, row in df.iterrows():
            client_scores = calculate_category_scores(row, QUESTION_CATEGORIES, CLEANED_TO_ORIGINAL_COL)
            client_scores['Email Address'] = row['Email Address']
            client_scores['Timestamp'] = row['Timestamp']
            category_scores_list.append(client_scores)
        
        scores_df = pd.DataFrame(category_scores_list)
        print(f"✓ Calculated scores for {len(scores_df)} clients")
        
        # Generate presentations
        print(f"\nGenerating presentations...")
        generated_files = []
        skipped_files = []

        for idx, row in scores_df.iterrows():
            client_email = row['Email Address']
            client_scores = {cat: row[cat] for cat in QUESTION_CATEGORIES.keys()}

            client_row = df[df['Email Address'] == client_email].iloc[0]
            client_responses = {}
            for col in question_columns:
                client_responses[col] = client_row[col]

            filename = email_to_filename(client_email)
            output_filename = os.path.join(OUTPUT_DIR, filename)
            
            # Check if presentation already exists
            if os.path.exists(output_filename):
                print(f"\n  Skipping: {client_email} (presentation already exists)")
                skipped_files.append(output_filename)
                continue
            
            print(f"\n  Processing: {client_email}")
            try:
                generate_client_presentation(
                    client_email, client_scores, client_responses,
                    COLUMN_NAME_MAPPING, CLEANED_TO_ORIGINAL_COL,
                    QUESTION_CATEGORIES, output_filename
                )
                generated_files.append(output_filename)
                print(f"    ✓ Saved: {filename}")
            except Exception as e:
                print(f"    ❌ Error: {e}")

        print(f"\n{'='*60}")
        print(f"✓ Successfully generated {len(generated_files)} new presentations")
        if skipped_files:
            print(f"⊘ Skipped {len(skipped_files)} existing presentations")
        print(f"Saved to: {OUTPUT_DIR}/")
        
        # Upload all files to SharePoint at the end (separate step)
        if SHAREPOINT_UPLOAD_ENABLED:
            print(f"\n{'='*60}")
            print("SharePoint Upload Process")
            print(f"{'='*60}")
            
            # Get all .pptx files in output directory (both new and existing)
            all_output_files = []
            if os.path.exists(OUTPUT_DIR):
                for file in os.listdir(OUTPUT_DIR):
                    if file.endswith('.pptx') and not file.startswith('~$'):  # Exclude temp files
                        file_path = os.path.join(OUTPUT_DIR, file)
                        all_output_files.append(file_path)
            
            if not all_output_files:
                print("  No presentation files found in output directory")
            else:
                print(f"  Found {len(all_output_files)} presentation file(s) to check/upload")
                uploaded_count = 0
                skipped_count = 0
                failed_count = 0
                
                for file_path in all_output_files:
                    filename = os.path.basename(file_path)
                    print(f"\n  Checking: {filename}")
                    
                    result = upload_to_sharepoint(file_path, filename)
                    if result:
                        # Check the message to see if it was uploaded or skipped
                        # (The function prints the status, so we track based on that)
                        uploaded_count += 1
                    else:
                        failed_count += 1
                
                print(f"\n  SharePoint Summary:")
                print(f"    ✓ Processed: {uploaded_count} files")
                if failed_count > 0:
                    print(f"    ❌ Failed: {failed_count} files")
        else:
            print(f"\n  SharePoint upload disabled (set SHAREPOINT_UPLOAD=true to enable)")
        
        print(f"\n{'='*60}")
        print(f"Completed at: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        print("="*60)
        
    except Exception as e:
        print(f"\n❌ Fatal error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
